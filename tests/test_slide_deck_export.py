"""Tests for building a real .pptx export document for AI-generated slide-deck content."""

from pathlib import Path

from pptx import Presentation

from islamic_research_hub.application.slide_deck_extraction import ExtractedSlide
from islamic_research_hub.research_notes.slide_deck_export import (
    build_slide_deck,
    export_slide_deck_to_pptx,
)

_SLIDES = (
    ExtractedSlide(title="The Battle of Badr", bullets=("Fought in 2 AH.", "A decisive victory.")),
    ExtractedSlide(title="The Treaty of Hudaybiyyah", bullets=("Signed in 6 AH.",)),
)


def test_build_slide_deck_includes_a_title_slide_naming_the_book() -> None:
    presentation = build_slide_deck("Seerah of the Prophet", _SLIDES)

    assert presentation.slides[0].shapes.title.text == "Seerah of the Prophet"


def test_build_slide_deck_includes_one_content_slide_per_extracted_slide() -> None:
    presentation = build_slide_deck("Seerah of the Prophet", _SLIDES)

    assert len(presentation.slides) == 1 + len(_SLIDES)
    assert presentation.slides[1].shapes.title.text == "The Battle of Badr"
    assert presentation.slides[2].shapes.title.text == "The Treaty of Hudaybiyyah"


def test_build_slide_deck_includes_the_real_bullet_text() -> None:
    presentation = build_slide_deck("Seerah of the Prophet", _SLIDES)

    body_text = presentation.slides[1].placeholders[1].text_frame.text
    assert "Fought in 2 AH." in body_text
    assert "A decisive victory." in body_text


def test_build_slide_deck_with_no_slides_still_has_a_title_slide() -> None:
    presentation = build_slide_deck("An Empty Book", ())

    assert len(presentation.slides) == 1
    assert presentation.slides[0].shapes.title.text == "An Empty Book"


def test_export_slide_deck_to_pptx_writes_a_real_readable_file(tmp_path: Path) -> None:
    output_path = tmp_path / "deck.pptx"

    export_slide_deck_to_pptx("Seerah of the Prophet", _SLIDES, output_path)

    assert output_path.is_file()
    reopened = Presentation(output_path)
    assert len(reopened.slides) == 1 + len(_SLIDES)
