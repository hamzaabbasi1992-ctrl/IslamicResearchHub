"""Export real, AI-generated slide-deck content (Phase 17 Milestone 1)
as a real .pptx presentation.

Lives alongside `docx_writer.py`/`collection_export.py`/`ai_answer_export.py`
(the other real python-docx/python-pptx-touching code in this project)
rather than under `application/` - all are real storage-adjacent
adapters, not pure domain logic.
"""

from pathlib import Path

from pptx import Presentation

from islamic_research_hub.application.slide_deck_extraction import ExtractedSlide

_TITLE_LAYOUT_INDEX = 0
_TITLE_AND_CONTENT_LAYOUT_INDEX = 1


def build_slide_deck(book_title: str, slides: tuple[ExtractedSlide, ...]) -> Presentation:
    """Return a real, in-memory `Presentation` - a title slide naming the
    real source book, then one real content slide per `ExtractedSlide`
    (heading + bullet points), in the same order `slides` was given.

    Returns the `Presentation` object rather than saving it, so this
    stays testable without touching a real file -
    `export_slide_deck_to_pptx()` below is the thin save wrapper real
    callers use.
    """
    presentation = Presentation()
    title_slide = presentation.slides.add_slide(presentation.slide_layouts[_TITLE_LAYOUT_INDEX])
    title_slide.shapes.title.text = book_title
    if len(title_slide.placeholders) > 1:
        title_slide.placeholders[1].text = "AI-generated, grounded in this library's real content"

    for slide in slides:
        content_slide = presentation.slides.add_slide(
            presentation.slide_layouts[_TITLE_AND_CONTENT_LAYOUT_INDEX]
        )
        content_slide.shapes.title.text = slide.title
        body = content_slide.placeholders[1].text_frame
        body.text = slide.bullets[0]
        for bullet in slide.bullets[1:]:
            paragraph = body.add_paragraph()
            paragraph.text = bullet
    return presentation


def export_slide_deck_to_pptx(
    book_title: str, slides: tuple[ExtractedSlide, ...], output_path: Path
) -> None:
    """Build and save a real slide deck export to `output_path`."""
    build_slide_deck(book_title, slides).save(output_path)
